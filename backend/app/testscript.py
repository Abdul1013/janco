from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define Corporate/Minimalist Colors
    # Dark Blue (Primary): 0, 51, 102
    # Light Grey (Background): 245, 245, 245
    # Accent (Teal): 0, 128, 128
    
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255) # White background for cleanliness

    def format_title(shape, text):
        shape.text = text
        p = shape.text_frame.paragraphs
        p.font.size = Pt(36)
        p.font.name = 'Helvetica'
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

    def format_body(shape, text_list):
        tf = shape.text_frame
        tf.clear()
        for item in text_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.name = 'Helvetica'
            p.space_after = Pt(14)
            p.level = 0

    # SLIDE 1: TITLE SLIDE
    slide_layout = prs.slide_layouts # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "A Software-Driven Service Automation\nfor On-Demand Cleaning Operations"
    subtitle.text = "The JANCO Project\n\nSubmitted By: Abdulrasheed Olabanji Abdulhaq\n(LCU/UG/22/23560)\n\nSupervisor: Miss Fatoki"
    
    # Customizing Title Font
    p = title.text_frame.paragraphs
    p.font.name = "Helvetica"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(0, 51, 102)

    # SLIDE 2: INTRODUCTION
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Introduction")
    
    content = format_body(slide.placeholders[1], content)

    # SLIDE 3: PROBLEM STATEMENT
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Problem Statement: The Trust & Efficiency Void")
    
    content = format_body(slide.placeholders[1], content)

    # SLIDE 4: AIM & OBJECTIVES
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Aim and Objectives")
    
    content = format_body(slide.placeholders[1], content)

    # SLIDE 5: METHODOLOGY - ARCHITECTURE
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Methodology: System Architecture")
    
    content =   format_body(slide.placeholders[1], content)

    # SLIDE 6: METHODOLOGY - ALGORITHMS
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Methodology: Key Algorithmic Modules")
    
    content =  format_body(slide.placeholders[1], content)

    # SLIDE 7: SIGNIFICANCE
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Significance of the Study")
    
    content = format_body(slide.placeholders[1], content)

    # SLIDE 8: EXPECTED RESULTS
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Expected Results")
    
    content =  format_body(slide.placeholders[1], content)

    # SLIDE 9: CONCLUSION
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Conclusion")
    
    content = format_body(slide.placeholders[1], content)

    # SLIDE 10: REFERENCES
    # SLIDE 10: REFERENCES
    slide = prs.slides.add_slide(slide_layout)
    format_title(slide.shapes.title, "Selected References")
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    references = ["Reference 1", "Reference 2", "Reference 3"]
    for item in references:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) # Smaller font for refs
        p.space_after = Pt(10)
    # Save the presentation
    prs.save('JANCO_Defense_Presentation.pptx')
    print("Presentation 'JANCO_Defense_Presentation.pptx' created successfully.")

if __name__ == "__main__":
    create_presentation()